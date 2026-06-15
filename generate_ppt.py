# -*- coding: utf-8 -*-
"""
生成 病历质控台账生成器 产品说明书 PPT
运行: python generate_ppt.py
"""
import os, sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Installing python-pptx...")
    import subprocess
    subprocess.run([sys.executable, '-m', 'pip', 'install', 'python-pptx', '--quiet'])
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

BASE = os.path.dirname(os.path.abspath(__file__))
OUTPUT = os.path.join(BASE, '病历质控台账生成器_产品说明书.pptx')

# Colors
GREEN_DARK = RGBColor(0x1A, 0x52, 0x76)
GREEN_MID = RGBColor(0x2E, 0x7D, 0x32)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x75, 0x75, 0x75)
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
ACCENT = RGBColor(0xC6, 0x28, 0x28)

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_title_bar(slide, title, subtitle=''):
    """Dark green top bar"""
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, Inches(1.2))  # 1=rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN_DARK
    shape.line.fill.background()
    add_text_box(slide, 0.8, 0.15, 11, 0.7, title, 32, WHITE, True)
    if subtitle:
        add_text_box(slide, 0.8, 0.75, 11, 0.4, subtitle, 14, RGBColor(0xBB, 0xCC, 0xDD))

def add_bullet_list(slide, left, top, width, height, items, font_size=16):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(6)
        p.level = item.count('  ') // 2
    return tf

def add_footer(slide, text='辽宁中医嘉和医院 · 病历质控台账生成器 v2.0 · 产品说明书'):
    add_text_box(slide, 0.5, 7.0, 12, 0.4, text, 10, GRAY, alignment=PP_ALIGN.CENTER)

# ==========================================
# Slide 1: Title
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, GREEN_DARK)
add_text_box(slide, 1, 1.5, 11, 1.2, '病历质控台账生成器', 48, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 2.8, 11, 0.8, '辽宁中医嘉和医院 · 院级质控专用', 24, RGBColor(0xBB, 0xCC, 0xDD), alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.8, 11, 0.6, 'v2.0  |  2026年6月', 18, RGBColor(0x99, 0xAA, 0xBB), alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.2, 11, 0.9, 'XPS提取 → AI六维分析 → 缺陷台账 → 一键导出Excel\n全流程自动化，人工仅需复核', 16, RGBColor(0xAA, 0xBB, 0xCC), alignment=PP_ALIGN.CENTER)

# ==========================================
# Slide 2: Pain Points → Solution
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '产品定位：让病历质控从"人查"到"机查"', '痛点 → 解决方案')

items = [
    '⏱️ 人工查1份病历 30-60 分钟 → AI 分析 15-30 秒，人工仅复核',
    '📏 质控标准执行不一致 → 六维 34 项标准内置，每次检查全覆盖',
    '📝 缺陷描述不规范 → AI 润色引擎，口语一键转专业质控语言',
    '📊 台账整理费时费力 → 自动去重合并、按人/按维度统计',
    '🔄 历史缺陷无法追溯 → 本地历史记录，常用描述一键复用',
    '🚫 知识产权无保护 → HTML 源码嵌入 EXE，F12 无法查看',
]
add_bullet_list(slide, 0.8, 1.6, 11.5, 5, items, 18)
add_footer(slide)

# ==========================================
# Slide 3: Six Dimensions
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '六维质控标准体系', '34项检查细则 + 5项一票否决')

dims = [
    ('一、核心时限', '7项', '入院记录24h · 首程8h · 主治48h · 主任72h\n连续3天病程 · 康复评定72h · 阶段小结30天'),
    ('二、康复专项', '4项', '标准化量表 · 中期评定4周 · 治疗记录单完整\n医嘱单/治疗单/病程记录三单一致'),
    ('三、病程记录质量', '6项', '首程要素完整 · 上级查房分析 · 病情变化记录\n有创操作即时记录 · 患方签名 · 阶段小结非复制'),
    ('四、文书完整性', '8项', '一般项目无空项 · 三史无遗漏 · 体格检查完整\n辅助检查有序 · 报告单签名 · 医师亲笔签名 · 首页完整'),
    ('五、知情同意与告知', '5项', '授权委托书 · 特殊检查同意书 · 康复风险告知书\n各类告知书齐全 · 患方签名+日期'),
    ('六、书写规范', '4项', '无大段复制粘贴 · 双线划改规范\n24小时制记录 · 术语规范无矛盾'),
]

for i, (dim_name, count, desc) in enumerate(dims):
    y = 1.5 + i * 0.95
    # Dimension label
    shape = slide.shapes.add_shape(
        1, Inches(0.6), Inches(y), Inches(2.8), Inches(0.8))
    shape.fill.solid()
    color = [GREEN_DARK, RGBColor(0x24,0x71,0xA3), RGBColor(0x2E,0x86,0xC1),
             RGBColor(0x34,0x98,0xDB), RGBColor(0x5D,0xAD,0xE2), RGBColor(0x85,0xC1,0xE9)][i]
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    add_text_box(slide, 0.8, y + 0.1, 2.5, 0.6, f'{dim_name}  [{count}]', 14, WHITE, True)
    add_text_box(slide, 3.6, y + 0.05, 9, 0.75, desc, 13, BLACK)

# Veto
shape = slide.shapes.add_shape(1, Inches(0.6), Inches(7.2), Inches(12), Inches(0.25))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()
add_text_box(slide, 0.8, 7.2, 11.5, 0.25, '🚨 一票否决：核心记录超时24h | 康复初始评定完全缺失 | 知情同意书缺失/伪造 | 诊断与治疗严重不符 | 复制粘贴致重大错误', 11, WHITE, True)

add_footer(slide, '')

# ==========================================
# Slide 4: Automation Pipeline
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '自动化流水线', '三步完成：XPS提取 → AI分析 → 导入台账')

steps = [
    ('📑', 'XPS 提取', 'HIS导出的 .xps → 一键提取 →\n纯文本病历 (_extracted/)', GREEN_DARK),
    ('🔬', 'AI 六维分析', '病历文本 → DeepSeek 逐项检查 →\n结构化缺陷 JSON (_analysis/)', GREEN_MID),
    ('📋', '导入复核', 'JSON → 自动填表 → 人工复核 →\n添加到台账 → 导出 Excel', RGBColor(0x8E, 0x44, 0xAD)),
]
for i, (icon, title, desc, color) in enumerate(steps):
    x = 1 + i * 4
    # Step box
    shape = slide.shapes.add_shape(1, Inches(x), Inches(2.0), Inches(3.5), Inches(4.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    add_text_box(slide, x + 0.2, 2.2, 3, 0.6, f'{icon}  Step {i+1}', 16, color, True)
    add_text_box(slide, x + 0.2, 2.8, 3, 0.6, title, 20, BLACK, True)
    add_text_box(slide, x + 0.2, 3.6, 3, 1.0, desc, 14, BLACK)
    if i < 2:
        add_text_box(slide, x + 3.5, 3.5, 0.5, 0.5, '→', 28, GRAY, True, PP_ALIGN.CENTER)

# Full pipeline bar
shape = slide.shapes.add_shape(1, Inches(1), Inches(6.5), Inches(11.3), Inches(0.5))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xF8, 0xF4, 0xFF)
shape.line.color.rgb = RGBColor(0xD5, 0xC8, 0xF0)
add_text_box(slide, 1.2, 6.5, 11, 0.5, '🚀 一键全流程：提取全部 XPS → 分析全部待处理 → 全部导入台账  （无需人工干预）', 16, RGBColor(0x8E, 0x44, 0xAD), True, PP_ALIGN.CENTER)
add_footer(slide)

# ==========================================
# Slide 5: Key Features
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '核心功能', '')

features = [
    ('🤖 AI 智能分析', 'DeepSeek 大模型逐份逐项检查\n自动识别患者信息、病区、主管医师\n仅输出有缺陷项，合格项不列\n一票否决自动标注'),
    ('✨ AI 文字润色', '口语描述 → 专业质控语言\n支持单条/全量一键润色\n历史描述本地存储，点击复用'),
    ('📥 Excel 导出', '绿色医院主题配色\nSheet 1: 缺陷明细（全字段）\nSheet 2: 汇总统计（按维度/按人）\n一票否决行自动标红加粗'),
    ('📦 台账管理', '同住院号自动合并去重\n单元格双击编辑，实时生效\n支持逐条删除/一键清空\n实时统计病历份数/缺陷总数'),
]

for i, (title, desc) in enumerate(features):
    x = 0.6 + (i % 2) * 6.2
    y = 1.5 + (i // 2) * 2.8
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(5.8), Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    add_text_box(slide, x + 0.3, y + 0.2, 5, 0.5, title, 20, GREEN_DARK, True)
    add_text_box(slide, x + 0.3, y + 0.8, 5, 1.5, desc, 14, BLACK)

add_footer(slide)

# ==========================================
# Slide 6: Usage
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '三种操作方式', '灵活适配不同场景')

methods = [
    ('方式一：自动化面板（推荐）', [
        '点击 🤖导入AI结果 → 打开自动化面板',
        '📑 XPS提取 → 🚀 一键分析 → 📊 导入选中',
        '每条缺陷自动匹配到六维表单对应位置',
        '人工逐条复核 → 修改 → 确认责任人',
        '点击 ➕添加到台账 → 📥导出全部',
    ]),
    ('方式二：全流程一键自动化', [
        '自动化面板底部点击 🚀启动全流程',
        '自动：提取XPS → 分析 → 导入台账',
        '全程无需人工干预',
        '完成后直接进入台账复核导出',
    ]),
    ('方式三：手动填写', [
        '直接在六维表单中逐项描述缺陷',
        '点击 ✨ 按钮润色单条描述',
        '勾选一票否决项',
        '添加到台账 → 导出 Excel',
    ]),
]

for i, (title, items) in enumerate(methods):
    x = 0.6 + i * 4.2
    shape = slide.shapes.add_shape(1, Inches(x), Inches(1.6), Inches(3.9), Inches(5.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFA, 0xFB, 0xFC)
    shape.line.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
    add_text_box(slide, x + 0.2, 1.8, 3.5, 0.5, title, 16, GREEN_DARK, True)
    for j, item in enumerate(items):
        add_text_box(slide, x + 0.3, 2.5 + j * 0.7, 3.3, 0.6, f'• {item}', 12, BLACK)

add_footer(slide)

# ==========================================
# Slide 7: Technical Architecture
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '技术架构与知识产权保护', '')

tech_items = [
    '后端框架：Python Flask + waitress 生产级 WSGI',
    'AI 引擎：DeepSeek Chat API (deepseek-chat)',
    '前端：单页 HTML5 + Vanilla JavaScript + ExcelJS',
    '打包：PyInstaller 单文件 EXE，双击即用',
    '平台：Windows 10/11 x64，无需安装 Python',
    '',
    '🔒 知识产权保护：',
    '• HTML 源码嵌入 EXE 二进制 → F12 无法查看',
    '• EXE 运行时无控制台窗口 → 后台静默运行',
    '• 自定义医院 LOGO 图标嵌入',
    '• 不写注册表 / 不创建系统服务 / 无残留文件',
    '',
    '📡 10 个 API 端点：',
    '/api/polish · /api/xps-list · /api/extract-xps · /api/extract-all',
    '/api/extracted-list · /api/run-analysis · /api/analysis-list · /api/analysis',
]
add_bullet_list(slide, 0.8, 1.5, 11.5, 5, tech_items, 16)
add_footer(slide)

# ==========================================
# Slide 8: Deployment
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '部署与运维', '开箱即用，零配置部署')

deploy = [
    '📂 目录结构：',
    '  病历质控台账生成器.exe          ← 主程序',
    '  config.json                     ← API Key 配置',
    '  data/病历质控测试组/            ← 病历数据',
    '    ├── A组/ C组/ D组/            ← XPS 源文件',
    '    ├── _extracted/                ← 提取文本',
    '    └── _analysis/                 ← AI 分析结果',
    '',
    '📋 新增病历：将 HIS 导出的 XPS 放入对应病区文件夹 → 重启程序',
    '💾 备份建议：定期备份 data/ 目录及导出的 Excel 台账',
    '📋 迁移部署：复制整个程序目录到目标电脑即可',
    '🔑 API Key：在网页顶部工具栏配置，浏览器本地存储，跨会话保留',
]
add_bullet_list(slide, 0.8, 1.5, 11.5, 5, deploy, 16)
add_footer(slide)

# ==========================================
# Slide 9: FAQ
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '常见问题', '')

faqs = [
    'Q: 双击 EXE 无反应？',
    'A: 检查 config.json，确认 8081 端口未被占用。',
    '',
    'Q: AI 分析返回 "LLM not configured"？',
    'A: 在页面顶部填入有效的 DeepSeek API Key (sk-xxx)。',
    '',
    'Q: 匹配数为 0？',
    'A: 确认 _analysis/ 下有 JSON 文件，点击刷新重新加载。',
    '',
    'Q: 如何迁移到其他电脑？',
    'A: 复制整个程序目录即可。可能需要安装 VC++ Redistributable。',
    '',
    'Q: 如何卸载？',
    'A: 直接删除程序目录，不写注册表，无残留文件。',
    '',
    'Q: AI 分析准确吗？',
    'A: AI 结果仅供参考，最终须由具有资质的质控人员复核确认。',
]
add_bullet_list(slide, 0.8, 1.5, 11.5, 5.5, faqs, 15)
add_footer(slide)

# ==========================================
# Slide 10: Thank You
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, GREEN_DARK)
add_text_box(slide, 1, 2.0, 11, 1.0, '谢谢！', 48, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.2, 11, 0.8, '病历质控台账生成器 v2.0', 28, RGBColor(0xBB, 0xCC, 0xDD), alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.3, 11, 1.2, '辽宁中医嘉和医院 · 医务科/医保办\n2026年6月', 18, RGBColor(0x99, 0xAA, 0xBB), alignment=PP_ALIGN.CENTER)

# Save
prs.save(OUTPUT)
print(f'PPT saved: {OUTPUT}')
